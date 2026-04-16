"""Generate SIGMA How-To PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x1a, 0x1a, 0x2e)   # deep navy
ACCENT  = RGBColor(0x4a, 0x9e, 0xff)   # sky blue
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xcc, 0xdd, 0xf0)
MUTED   = RGBColor(0x88, 0x99, 0xbb)
YELLOW  = RGBColor(0xff, 0xd7, 0x66)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, left, top, width, height,
          size=Pt(18), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, top=Inches(0.08)):
    """Thin horizontal accent line near top."""
    box(slide, Inches(0.5), top, SLIDE_W - Inches(1), Pt(3),
        fill_color=ACCENT)


def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    label(slide, title,
          Inches(0.5), Inches(0.22), SLIDE_W - Inches(1), Inches(0.7),
          size=Pt(28), bold=True, color=ACCENT)
    if subtitle:
        label(slide, subtitle,
              Inches(0.5), Inches(0.9), SLIDE_W - Inches(1), Inches(0.45),
              size=Pt(14), color=MUTED)


def bullet_section(slide, left, top, width, height,
                   heading, items, head_color=YELLOW, item_color=LIGHT,
                   item_size=Pt(14), head_size=Pt(15)):
    """Render a headed bullet list inside a rounded box."""
    shape = box(slide, left, top, width, height,
                line_color=ACCENT, line_width=Pt(1.2))
    # heading
    label(slide, heading,
          left + Inches(0.15), top + Inches(0.1),
          width - Inches(0.3), Inches(0.35),
          size=head_size, bold=True, color=head_color)
    # items
    y = top + Inches(0.45)
    for item in items:
        label(slide, f"• {item}",
              left + Inches(0.2), y,
              width - Inches(0.4), Inches(0.38),
              size=item_size, color=item_color)
        y += Inches(0.34)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# large centre title
label(s, "Σ",
      Inches(5.4), Inches(0.9), Inches(2.5), Inches(1.1),
      size=Pt(72), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

label(s, "SIGMA",
      Inches(1), Inches(1.9), Inches(11.33), Inches(1.1),
      size=Pt(72), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

label(s, "Segmentation & Image Guided Medical Annotation",
      Inches(1.5), Inches(3.0), Inches(10.33), Inches(0.7),
      size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER)

label(s, "View and segment DICOM and NIfTI volumes in your browser — no desktop install required.",
      Inches(2), Inches(3.75), Inches(9.33), Inches(0.5),
      size=Pt(15), color=MUTED, align=PP_ALIGN.CENTER)

box(s, Inches(4.5), Inches(4.55), Inches(4.33), Pt(3), fill_color=ACCENT)

label(s, "How-To Guide",
      Inches(1), Inches(4.75), Inches(11.33), Inches(0.5),
      size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 2 – What is SIGMA?
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "What is ΣIGMA?",
            "A browser-native ITK-SNAP alternative for researchers and radiologists")

cols = [
    ("Viewer", [
        "Axial, Coronal, Sagittal + Oblique",
        "Synchronized crosshairs",
        "Single-panel zoom mode",
        "Window / Level presets",
    ]),
    ("Editor", [
        "Paint & erase segmentation masks",
        "Region grow (Grow2D)",
        "Fill holes, clear slice",
        "Undo (3 levels)",
    ]),
    ("Workflow", [
        "Folder-based volume catalog",
        "DICOM & NIfTI support",
        "Save masks as NIfTI",
        "AI segmentation integration",
    ]),
]

x = Inches(0.4)
for heading, items in cols:
    bullet_section(s, x, Inches(1.5), Inches(4.1), Inches(5.3),
                   heading, items)
    x += Inches(4.3)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 3 – Getting Started
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Getting Started", "Install once, then just run start.sh")

steps = [
    ("1", "Clone the repo",
     "git clone https://github.com/slowvak/SIGMA.git"),
    ("2", "Set up the Python environment",
     "cd server && uv sync"),
    ("3", "Install frontend dependencies",
     "cd client && npm install"),
    ("4", "Launch both servers",
     "./start.sh   →   open the URL shown in the terminal"),
]

y = Inches(1.55)
for num, title_text, detail in steps:
    # circle number
    circ = box(s, Inches(0.35), y + Inches(0.05),
               Inches(0.55), Inches(0.55),
               fill_color=ACCENT)
    label(s, num,
          Inches(0.35), y + Inches(0.05), Inches(0.55), Inches(0.55),
          size=Pt(18), bold=True, color=BG, align=PP_ALIGN.CENTER)
    label(s, title_text,
          Inches(1.1), y, Inches(11), Inches(0.38),
          size=Pt(16), bold=True, color=WHITE)
    label(s, detail,
          Inches(1.1), y + Inches(0.35), Inches(11), Inches(0.35),
          size=Pt(13), color=MUTED)
    y += Inches(1.1)

label(s, "Prerequisites: uv (Python package manager) and Node.js + npm",
      Inches(0.4), Inches(6.9), Inches(12), Inches(0.4),
      size=Pt(12), color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 4 – Opening a Volume
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Opening a Volume")

flow = [
    ("Click  Open Folder",
     "Select any directory on your filesystem — ΣIGMA scans it automatically for DICOM series and NIfTI files."),
    ("Pick a volume from the list",
     "The left panel shows every discovered series. Click one to load it into the viewer."),
    ("Wait for the progress bar",
     "Volume data streams from the local server into browser memory. Large CTs take a few seconds."),
    ("Use  Back to Volumes  to return",
     "Switch volumes at any time without restarting the server."),
]

y = Inches(1.5)
for i, (step, desc) in enumerate(flow):
    box(s, Inches(0.4), y, Inches(12.5), Inches(0.92),
        line_color=ACCENT, line_width=Pt(0.8))
    label(s, step,
          Inches(0.65), y + Inches(0.06), Inches(5.5), Inches(0.38),
          size=Pt(15), bold=True, color=YELLOW)
    label(s, desc,
          Inches(0.65), y + Inches(0.44), Inches(12), Inches(0.38),
          size=Pt(13), color=LIGHT)
    y += Inches(1.05)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 5 – The 4-Panel Viewer
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "The 4-Panel Viewer", "Axial · Coronal · Sagittal · Oblique — all synchronized")

# draw a representative 2×2 grid
GX, GY = Inches(0.4), Inches(1.35)
GW, GH = Inches(5.6), Inches(5.7)
CELL_W = GW / 2
CELL_H = GH / 2
labels_grid = ["Axial", "Coronal", "Sagittal", "Oblique"]
positions = [(GX, GY), (GX + CELL_W, GY),
             (GX, GY + CELL_H), (GX + CELL_W, GY + CELL_H)]

for (lx, ly), lbl in zip(positions, labels_grid):
    box(s, lx + Inches(0.04), ly + Inches(0.04),
        CELL_W - Inches(0.08), CELL_H - Inches(0.08),
        line_color=ACCENT, line_width=Pt(1.5))
    label(s, lbl,
          lx + Inches(0.15), ly + Inches(0.12), CELL_W, Inches(0.35),
          size=Pt(13), bold=True, color=ACCENT)

# callouts on the right
RX = Inches(6.4)
tips = [
    ("Scroll slices",
     "Mouse wheel or click-drag up/down on any panel"),
    ("Crosshair sync",
     "All panels update instantly when you click with the crosshair tool"),
    ("Single-panel mode",
     "Click the panel's name label (A / C / S) to expand it full-screen; press 4 to restore"),
    ("Slice slider",
     "The bar below each panel lets you jump to any slice position"),
    ("Oblique panel",
     "Always shows the plane defined by the current crosshair intersection"),
]
y = Inches(1.35)
for head, body in tips:
    label(s, head, RX, y, Inches(6.5), Inches(0.3),
          size=Pt(14), bold=True, color=YELLOW)
    label(s, body, RX + Inches(0.15), y + Inches(0.28), Inches(6.3), Inches(0.4),
          size=Pt(12), color=LIGHT)
    y += Inches(0.88)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 6 – Window / Level
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Window / Level",
            "Adjust image contrast and brightness for your tissue of interest")

presets = [
    ("Brain",  "W 80  L 40"),
    ("Bone",   "W 2000  L 400"),
    ("Lung",   "W 1500  L -600"),
    ("Abdomen","W 350  L 40"),
]

label(s, "Right-click drag on any viewer panel:",
      Inches(0.5), Inches(1.4), Inches(12), Inches(0.4),
      size=Pt(16), bold=True, color=WHITE)

arrows = [
    ("← →  (left / right)", "Changes Window width (contrast)"),
    ("↑ ↓  (up / down)",    "Changes Level / centre point (brightness)"),
]
y = Inches(1.9)
for key, desc in arrows:
    label(s, key,  Inches(0.8), y, Inches(4.5), Inches(0.38),
          size=Pt(15), bold=True, color=ACCENT)
    label(s, desc, Inches(5.4), y, Inches(7.5), Inches(0.38),
          size=Pt(14), color=LIGHT)
    y += Inches(0.52)

label(s, "Quick Presets",
      Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
      size=Pt(16), bold=True, color=WHITE)

px = Inches(0.5)
for name, vals in presets:
    b = box(s, px, Inches(3.6), Inches(2.8), Inches(1.1),
            fill_color=RGBColor(0x1e, 0x2e, 0x4a), line_color=ACCENT)
    label(s, name,
          px, Inches(3.68), Inches(2.8), Inches(0.4),
          size=Pt(15), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    label(s, vals,
          px, Inches(4.05), Inches(2.8), Inches(0.35),
          size=Pt(12), color=LIGHT, align=PP_ALIGN.CENTER)
    px += Inches(3.1)

label(s, "Current W/L values are shown in the tool panel at all times.",
      Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
      size=Pt(13), color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 7 – Segmentation Tools
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Segmentation Tools")

tools = [
    ("✛  Crosshair",
     ["Navigate across all panels simultaneously",
      "Click to set crosshair; all views re-centre"]),
    ("✏  Paint",
     ["Left-click drag to paint the active label",
      "Brush Radius slider controls size",
      "Brush Depth paints through N adjacent slices",
      "Intensity Limits restrict painting by HU range"]),
    ("⊕  Grow2D",
     ["Click a seed voxel on the current slice",
      "Expands to connected voxels within Min–Max intensity",
      "Range auto-sets to mean±stdev of 5×5 patch around seed",
      "Adjust with the dual slider or type values directly"]),
    ("⊘  Erase",
     ["Select Erase mode in the label panel",
      "Brush removes voxels of any label"]),
]

y = Inches(1.4)
for i, (name, pts) in enumerate(tools):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * Inches(6.5)
    cy = y + row * Inches(2.65)
    h = Inches(0.38 + len(pts) * 0.34 + 0.25)
    box(s, cx, cy, Inches(6.1), h,
        line_color=ACCENT, line_width=Pt(1))
    label(s, name, cx + Inches(0.15), cy + Inches(0.1),
          Inches(5.8), Inches(0.38),
          size=Pt(15), bold=True, color=YELLOW)
    iy = cy + Inches(0.48)
    for pt in pts:
        label(s, f"• {pt}", cx + Inches(0.25), iy,
              Inches(5.6), Inches(0.34),
              size=Pt(12), color=LIGHT)
        iy += Inches(0.33)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 8 – Actions
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Actions", "Operations in the toolbar that modify the segmentation")

actions = [
    ("Undo  (Ctrl+Z)",
     "Reverts the last paint, grow, refine, propagate, or fill operation. Supports up to 3 levels of undo."),
    ("Refine",
     "Snaps the active label boundary to image edges on the current axial slice using Sobel gradient detection."),
    ("Propagate",
     "Copies the label from the adjacent slice and refines it — use this to step through a stack slice by slice."),
    ("Fill Holes",
     "Fills enclosed background regions within each connected component of the active label on the current slice."),
    ("Clear Slice",
     "Removes all voxels of the active label on the current slice only. Useful for restarting a single slice."),
    ("Filter",
     "Smooths raw image intensities. Choose 2D/3D, Mean/Median/Sigma (Gaussian), kernel size 3/5/7, Slice or Volume."),
    ("Save As…",
     "Writes the current segmentation back to the server as a NIfTI (.nii.gz) mask file."),
]

y = Inches(1.4)
col_w = Inches(6.1)
for i, (name, desc) in enumerate(actions):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * Inches(6.6)
    cy = y + row * Inches(1.2)
    if i == len(actions) - 1 and len(actions) % 2 == 1:
        cx = Inches(3.7)  # centre the lone last item
    box(s, cx, cy, col_w, Inches(1.05),
        line_color=ACCENT, line_width=Pt(0.8))
    label(s, name, cx + Inches(0.15), cy + Inches(0.08),
          col_w - Inches(0.3), Inches(0.35),
          size=Pt(14), bold=True, color=YELLOW)
    label(s, desc, cx + Inches(0.15), cy + Inches(0.44),
          col_w - Inches(0.3), Inches(0.5),
          size=Pt(11), color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 9 – Managing Labels
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Managing Labels", "The label panel lives on the right side of the toolbar")

items = [
    ("Single-click a label",
     "Makes it the active label for painting, erasing, and all segmentation operations."),
    ("Double-click a label",
     "Opens the label editor — rename the label, change its colour, or delete it."),
    ("Eye icon",
     "Toggle the visibility of a label in the overlay without changing which label is active."),
    ("Colour swatch",
     "Quick-click to change a label's display colour directly from the list."),
    ("+ button",
     "Add a new label. Prompts for a name and auto-assigns the next available colour."),
    ("Label Overlay Opacity",
     "Slider (0–100%) controls how opaque the segmentation colour overlay appears over the image."),
]

y = Inches(1.45)
for name, desc in items:
    label(s, name, Inches(0.5), y, Inches(5.0), Inches(0.38),
          size=Pt(14), bold=True, color=YELLOW)
    label(s, desc, Inches(5.6), y, Inches(7.4), Inches(0.55),
          size=Pt(13), color=LIGHT, wrap=True)
    box(s, Inches(0.5), y + Inches(0.44), Inches(12.3), Pt(0.8),
        fill_color=RGBColor(0x2a, 0x3a, 0x5a))
    y += Inches(0.82)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 10 – AI Integration
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "AI Integration", "Accelerate segmentation with automated models")

label(s, "Click the 🤖 AI button in the toolbar to open the model picker.",
      Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.45),
      size=Pt(15), color=WHITE)

options = [
    ("TotalSegmentator",
     "Downloads the current volume as a NIfTI file and opens totalsegmentator.com for full-body auto-segmentation. The result can be imported back as a label file.",
     ACCENT),
    ("Server-Side Models",
     "Custom models defined in  models/ai-models.json  appear in the picker automatically. Any model that accepts a NIfTI volume and returns a mask can be wired in.",
     YELLOW),
]

y = Inches(2.1)
for name, desc, col in options:
    b = box(s, Inches(0.5), y, Inches(12.3), Inches(1.8),
            line_color=col, line_width=Pt(1.5))
    label(s, name, Inches(0.75), y + Inches(0.15), Inches(11.5), Inches(0.45),
          size=Pt(17), bold=True, color=col)
    label(s, desc, Inches(0.75), y + Inches(0.6), Inches(11.5), Inches(1.0),
          size=Pt(13), color=LIGHT)
    y += Inches(2.05)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 11 – Keyboard Shortcuts
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Keyboard Shortcuts")

shortcuts = [
    ("?",          "Open the in-app help panel"),
    ("Ctrl + Z",   "Undo last segmentation edit (up to 3 levels)"),
    ("Escape",     "Close any open modal dialog"),
    ("4",          "Return to 4-panel view from single-panel mode"),
    ("A / C / S",  "Click the panel label to expand Axial / Coronal / Sagittal to full screen"),
    ("Mouse Wheel","Scroll through slices on the focused panel"),
    ("Right-click drag", "Adjust Window (left/right) and Level (up/down)"),
]

# two-column layout
mid = (len(shortcuts) + 1) // 2
left_col  = shortcuts[:mid]
right_col = shortcuts[mid:]

def draw_shortcuts(slide, items, cx, cy):
    for key, desc in items:
        b = box(slide, cx, cy, Inches(3.0), Inches(0.52),
                fill_color=RGBColor(0x1e, 0x2e, 0x4a), line_color=ACCENT, line_width=Pt(0.8))
        label(slide, key,
              cx + Inches(0.12), cy + Inches(0.08), Inches(2.8), Inches(0.38),
              size=Pt(14), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
        label(slide, desc,
              cx + Inches(3.2), cy + Inches(0.1), Inches(3.1), Inches(0.38),
              size=Pt(13), color=LIGHT)
        cy += Inches(0.7)

draw_shortcuts(s, left_col,  Inches(0.4),  Inches(1.45))
draw_shortcuts(s, right_col, Inches(6.75), Inches(1.45))


# ═══════════════════════════════════════════════════════════════════════════
# Slide 12 – Typical Workflow
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Typical Segmentation Workflow")

steps = [
    ("Open Folder",           "Point ΣIGMA at the directory containing your DICOMs or NIfTIs"),
    ("Select Volume",         "Click the series in the left panel to load it"),
    ("Set Window / Level",    "Use a preset (Brain / Bone / Lung / Abd) or right-click drag to tune contrast"),
    ("Create / select label", "Click + to add a label, or single-click an existing one to activate it"),
    ("Paint seed",            "Use Paint or Grow2D to lay down your initial segmentation on a representative slice"),
    ("Propagate",             "Press Propagate to copy-and-refine the label to adjacent slices, stepping through the stack"),
    ("Clean up",              "Use Fill Holes, Clear Slice, or manual painting to fix any errors"),
    ("Save",                  "Click Save As… to write the mask back to disk as a NIfTI file"),
]

ARROW_COLOR = ACCENT
y = Inches(1.4)
for i, (step, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * Inches(6.5)
    cy = y + row * Inches(1.35)
    box(s, cx, cy, Inches(6.0), Inches(1.15),
        line_color=ACCENT, line_width=Pt(0.8))
    # step number badge
    label(s, str(i + 1),
          cx + Inches(0.1), cy + Inches(0.1), Inches(0.5), Inches(0.5),
          size=Pt(20), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    label(s, step,
          cx + Inches(0.7), cy + Inches(0.08), Inches(5.0), Inches(0.38),
          size=Pt(14), bold=True, color=YELLOW)
    label(s, desc,
          cx + Inches(0.7), cy + Inches(0.48), Inches(5.1), Inches(0.5),
          size=Pt(11), color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 13 – Tips & Tricks
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Tips & Tricks")

tips_list = [
    ("Use Intensity Limits while painting",
     "Restrict the brush to a HU range so you only paint the tissue you intend — prevents spill-over into adjacent structures."),
    ("Grow2D + Propagate combo",
     "Grow2D a seed on the clearest slice, then Propagate through the rest of the stack for fast, consistent labelling."),
    ("Brush Depth for thick structures",
     "Set Brush Depth > 1 to paint through multiple slices at once — great for large uniform regions."),
    ("Adjust opacity while reviewing",
     "Lower the Label Overlay Opacity to see the underlying image anatomy while checking your mask boundaries."),
    ("Double-click to rename",
     "Double-click any label in the panel to rename it, change colour, or delete — keep labels descriptive from the start."),
    ("Filter before growing",
     "Apply a 2D Gaussian (Sigma) filter to noisy slices before using Grow2D — the smoother intensity surface yields cleaner region boundaries."),
]

y = Inches(1.4)
for i, (tip, body) in enumerate(tips_list):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * Inches(6.5)
    cy = y + row * Inches(1.7)
    box(s, cx, cy, Inches(6.1), Inches(1.55),
        line_color=ACCENT, line_width=Pt(0.8))
    label(s, tip,
          cx + Inches(0.15), cy + Inches(0.1), Inches(5.8), Inches(0.38),
          size=Pt(13), bold=True, color=YELLOW)
    label(s, body,
          cx + Inches(0.15), cy + Inches(0.48), Inches(5.8), Inches(0.95),
          size=Pt(11), color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out = "/Users/bje/repos/NextEd/SIGMA_HowTo.pptx"
prs.save(out)
print(f"Saved → {out}")
